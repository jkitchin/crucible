"""Source ingestion logic for Crucible.

Handles copying sources to the right location, extracting text,
and registering them in the database.
"""

import os
import re
import shutil
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path
from urllib.error import URLError
from urllib.request import Request, urlopen

from filelock import FileLock

from crucible.database import CrucibleDB
from crucible.orgparse import parse_org_file


def detect_source_type(path: Path) -> str:
    """Detect source type from file extension."""
    suffix = path.suffix.lower()
    type_map = {
        ".pdf": "pdf",
        ".org": "notebook",
        ".md": "web",
        ".html": "web",
        ".htm": "web",
        ".docx": "web",
        ".pptx": "web",
        ".csv": "data",
        ".json": "data",
        ".xlsx": "data",
        ".xls": "data",
        ".tsv": "data",
        ".hdf5": "data",
        ".h5": "data",
        ".npy": "data",
        ".npz": "data",
    }
    return type_map.get(suffix, "other")


CRUCIBLE_DIR = ".crucible"


def destination_dir(root: Path, source_type: str, shareable: bool) -> Path:
    """Determine where to store a source file."""
    base = root / CRUCIBLE_DIR / "sources"
    if shareable:
        if source_type == "notebook":
            return base / "notebooks"
        elif source_type == "data":
            return base / "data"
        else:
            return base / "data"
    else:
        if source_type == "pdf":
            return base / "external" / "pdfs"
        elif source_type == "web":
            return base / "external" / "web"
        else:
            return base / "external" / "web"


def atomic_write_text(path: Path, content: str, encoding: str = "utf-8"):
    """Write content to path atomically via temp file + rename."""
    fd, tmp_path = tempfile.mkstemp(dir=str(path.parent), suffix=".tmp")
    try:
        os.write(fd, content.encode(encoding))
        os.close(fd)
        fd = -1
        os.replace(tmp_path, str(path))
    except Exception:
        if fd >= 0:
            os.close(fd)
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
        raise


def extract_text(path: Path) -> str:
    """Extract plain text from a source file."""
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(path)
    elif suffix in (".html", ".htm"):
        return _extract_via_pandoc(path, input_format="html")
    elif suffix == ".docx":
        return _extract_via_pandoc(path, input_format="docx")
    elif suffix == ".pptx":
        return _extract_pptx(path)
    elif suffix in (".xlsx", ".xls"):
        return _extract_excel(path)
    elif suffix in (".org", ".md", ".txt"):
        return path.read_text(encoding="utf-8", errors="replace")
    elif suffix in (".csv", ".tsv"):
        return path.read_text(encoding="utf-8", errors="replace")
    else:
        return ""


def _extract_via_pandoc(path: Path, input_format: str) -> str:
    """Extract text via pandoc for formats it supports."""
    if not shutil.which("pandoc"):
        return path.read_text(encoding="utf-8", errors="replace") if input_format == "html" else ""
    result = subprocess.run(
        ["pandoc", "-f", input_format, "-t", "plain", "--wrap=none", str(path)],
        capture_output=True, text=True, timeout=60,
    )
    if result.returncode == 0:
        return result.stdout
    return ""


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using pdftotext or pandoc fallback."""
    if shutil.which("pdftotext"):
        result = subprocess.run(
            ["pdftotext", "-layout", str(path), "-"],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            return result.stdout
    return _extract_via_pandoc(path, input_format="pdf")


def _extract_pptx(path: Path) -> str:
    """Extract text from PowerPoint files."""
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    prs = Presentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append(" | ".join(cells))
        lines.append("")
    return "\n".join(lines)


def _extract_excel(path: Path) -> str:
    """Extract text from Excel files."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ""
    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append(" | ".join(cells))
        lines.append("")
    wb.close()
    return "\n".join(lines)


def extract_title(path: Path, text: str) -> str:
    """Try to extract a title from the source content."""
    suffix = path.suffix.lower()

    if suffix == ".org":
        meta = parse_org_file(path)
        if meta.title:
            return meta.title

    if suffix == ".md":
        for line in text.split("\n"):
            line = line.strip()
            if line.startswith("# "):
                return line[2:].strip()

    if suffix in (".html", ".htm"):
        import re
        m = re.search(r"<title>(.*?)</title>",
                       path.read_text(encoding="utf-8", errors="replace"),
                       re.IGNORECASE | re.DOTALL)
        if m:
            return m.group(1).strip()

    # Fallback to filename
    return path.stem.replace("-", " ").replace("_", " ").title()


def generate_cite_key(title: str, authors: list[str] | None, date: str | None) -> str:
    """Generate a BibTeX cite key from source metadata.

    Format: firstauthor_surname + year, e.g. smith2024.
    Falls back to title slug + year.
    """
    import re as _re
    year = (date or "")[:4] or "nd"

    if authors and authors[0]:
        # Take last name of first author
        first = authors[0]
        # Handle "Last, First" or "First Last"
        if "," in first:
            surname = first.split(",")[0].strip()
        else:
            surname = first.strip().split()[-1]
        surname = _re.sub(r"[^a-zA-Z]", "", surname).lower()
        return f"{surname}{year}"

    # Fallback to first significant word of title
    words = _re.sub(r"[^a-zA-Z\s]", "", title.lower()).split()
    stop = {"the", "a", "an", "of", "on", "in", "for", "and", "to", "with"}
    word = next((w for w in words if w not in stop), words[0] if words else "unknown")
    return f"{word}{year}"


def generate_bib_entry(cite_key: str, title: str, authors: list[str] | None,
                       date: str | None, url: str | None,
                       source_type: str, doi: str | None = None) -> str:
    """Generate a minimal BibTeX entry string.

    Used as a fallback when no DOI lookup or user-provided bibtex is available.
    Contains only the bibliographic fields we know — no crucible-internal metadata.
    """
    year = (date or "")[:4] or ""

    entry_type = {
        "pdf": "article",
        "web": "misc",
        "notebook": "misc",
        "data": "misc",
    }.get(source_type, "misc")

    lines = [f"@{entry_type}{{{cite_key},"]
    lines.append(f"  title = {{{title}}},")
    if authors:
        lines.append(f"  author = {{{' and '.join(authors)}}},")
    if year:
        lines.append(f"  year = {{{year}}},")
    if doi:
        lines.append(f"  doi = {{{doi}}},")
    if url:
        lines.append(f"  url = {{{url}}},")
    lines.append("}")
    return "\n".join(lines)


_CITE_KEY_RE = re.compile(r"(@\w+\s*\{\s*)([^,\s]+)")


def replace_cite_key(bibtex: str, new_key: str) -> str:
    """Replace the citation key of the first @entry in a bibtex string."""
    return _CITE_KEY_RE.sub(lambda m: m.group(1) + new_key, bibtex, count=1)


_BIB_ENTRY_RE = re.compile(r"@(\w+)\s*\{\s*([^,\s]+)\s*,(.*?)\n\}", re.DOTALL)
_BIB_FIELD_RE = re.compile(
    r"(\w+)\s*=\s*(?:\{((?:[^{}]|\{[^{}]*\})*)\}|\"([^\"]*)\")",
    re.DOTALL,
)


def _parse_bib_fields(body: str) -> dict[str, str]:
    """Parse the field body of a BibTeX entry into a dict."""
    fields: dict[str, str] = {}
    for m in _BIB_FIELD_RE.finditer(body):
        name = m.group(1).lower()
        value = m.group(2) if m.group(2) is not None else m.group(3) or ""
        fields[name] = " ".join(value.split())
    return fields


def parse_bib_file(bib_path: Path) -> list[dict]:
    """Parse a BibTeX file into a list of entries.

    Each entry dict has 'cite_key', 'entry_type', and parsed fields
    (title, author, year, doi, url, file, ...).
    """
    if not bib_path.exists():
        return []
    text = bib_path.read_text(encoding="utf-8", errors="replace")
    entries = []
    for m in _BIB_ENTRY_RE.finditer(text):
        entry = {
            "entry_type": m.group(1).lower(),
            "cite_key": m.group(2),
        }
        entry.update(_parse_bib_fields(m.group(3)))
        entries.append(entry)
    return entries


def _authors_from_bib(author_field: str) -> list[str]:
    """Split a BibTeX author field on ' and '."""
    if not author_field:
        return []
    return [a.strip() for a in re.split(r"\s+and\s+", author_field) if a.strip()]


def _source_type_for_entry(entry_type: str, path: Path | None) -> str:
    """Best-effort source_type for a bibtex entry."""
    if path is not None:
        t = detect_source_type(path)
        if t != "other":
            return t
    return {
        "article": "pdf",
        "inproceedings": "pdf",
        "book": "pdf",
        "phdthesis": "pdf",
        "mastersthesis": "pdf",
    }.get(entry_type, "other")


def upsert_sources_from_disk(root: Path, db) -> int:
    """Create sources rows for files in sources/notebooks/ and sources/data/.

    These directories are committed to git (unlike sources/external/), so a
    fresh clone has the files but no sources table. Scans both directories
    and creates a source row for each file, using the filename stem as the
    cite_key in metadata. Skips hidden files, extracted text sidecars
    (*.txt next to a non-txt source), and entries whose path already exists
    in the sources table.

    Returns the number of new sources inserted.
    """
    existing_paths: set[str] = {src["path"] for src in db.list_sources()}
    inserted = 0

    subdirs = [
        (root / CRUCIBLE_DIR / "sources" / "notebooks", "notebook"),
        (root / CRUCIBLE_DIR / "sources" / "data", "data"),
    ]

    for subdir, default_type in subdirs:
        if not subdir.is_dir():
            continue
        for path in sorted(subdir.rglob("*")):
            if not path.is_file() or path.name.startswith("."):
                continue
            # Skip *.txt sidecars from text extraction if a non-txt source
            # with the same stem exists
            if path.suffix.lower() == ".txt":
                siblings = [p for p in path.parent.glob(path.stem + ".*")
                            if p != path and p.is_file()]
                if siblings:
                    continue

            try:
                rel_path = str(path.relative_to(root.resolve()))
            except ValueError:
                rel_path = str(path)

            if rel_path in existing_paths:
                continue

            cite_key = path.stem
            title = path.stem.replace("-", " ").replace("_", " ").title()
            if path.suffix.lower() == ".org":
                try:
                    meta = parse_org_file(path)
                    if meta.title:
                        title = meta.title
                except Exception:
                    pass

            detected = detect_source_type(path)
            source_type = detected if detected != "other" else default_type

            metadata = {"cite_key": cite_key}

            try:
                mtime_iso = datetime.fromtimestamp(
                    path.stat().st_mtime
                ).strftime("%Y-%m-%d")
            except OSError:
                mtime_iso = None

            try:
                db.add_source(
                    path=rel_path,
                    title=title,
                    source_type=source_type,
                    shareable=True,
                    url=None,
                    authors=None,
                    date=mtime_iso,
                    metadata=metadata,
                )
                existing_paths.add(rel_path)
                inserted += 1
            except Exception:
                continue

    return inserted


def upsert_sources_from_bib(root: Path, db) -> int:
    """Create sources rows for every entry in .crucible/references.bib.

    Safe to call repeatedly: entries whose cite_key already exists in
    sources.metadata are skipped. Used by auto-rebuild (no DB) and by
    `crucible sync` so bib edits are picked up without re-ingesting.

    Returns the number of new sources inserted.
    """
    import json as _json

    bib_path = root / CRUCIBLE_DIR / "references.bib"
    entries = parse_bib_file(bib_path)
    if not entries:
        return 0

    existing_cite_keys: set[str] = set()
    for src in db.list_sources():
        raw = src.get("metadata") or "{}"
        try:
            meta = _json.loads(raw) if isinstance(raw, str) else (raw or {})
        except (ValueError, TypeError):
            meta = {}
        key = meta.get("cite_key")
        if key:
            existing_cite_keys.add(key)

    existing_paths: set[str] = {src["path"] for src in db.list_sources()}

    inserted = 0
    for entry in entries:
        cite_key = entry["cite_key"]
        if cite_key in existing_cite_keys:
            continue

        file_field = entry.get("file")
        rel_path: str | None = None
        abs_path: Path | None = None
        if file_field:
            candidate = Path(file_field)
            if not candidate.is_absolute():
                abs_path = (root / candidate).resolve()
            else:
                abs_path = candidate
            try:
                rel_path = str(abs_path.relative_to(root.resolve()))
            except ValueError:
                rel_path = str(abs_path)

        if rel_path is None:
            rel_path = f"{CRUCIBLE_DIR}/bib/{cite_key}"

        if rel_path in existing_paths:
            continue

        source_type = _source_type_for_entry(entry["entry_type"], abs_path)
        shareable = not rel_path.startswith(f"{CRUCIBLE_DIR}/sources/external/")
        authors = _authors_from_bib(entry.get("author", ""))
        date = entry.get("year") or None

        metadata: dict = {"cite_key": cite_key}
        if entry.get("doi"):
            metadata["doi"] = entry["doi"]

        try:
            db.add_source(
                path=rel_path,
                title=entry.get("title", cite_key),
                source_type=source_type,
                shareable=shareable,
                url=entry.get("url"),
                authors=authors,
                date=date,
                metadata=metadata,
            )
            existing_cite_keys.add(cite_key)
            existing_paths.add(rel_path)
            inserted += 1
        except Exception:
            continue

    return inserted


def fetch_bibtex_from_doi(doi: str, timeout: float = 10.0) -> str | None:
    """Fetch a full BibTeX entry for a DOI via doi.org content negotiation.

    Returns the entry string on success, or None if the lookup fails.
    """
    doi = doi.strip()
    if doi.lower().startswith("doi:"):
        doi = doi[4:]
    if doi.startswith("http"):
        url = doi
    else:
        url = f"https://doi.org/{doi}"

    req = Request(url, headers={
        "Accept": "application/x-bibtex; charset=utf-8",
        "User-Agent": "crucible (+https://github.com/jkitchin/crucible)",
    })
    try:
        with urlopen(req, timeout=timeout) as resp:
            data = resp.read().decode("utf-8", errors="replace").strip()
    except (URLError, OSError, TimeoutError):
        return None
    if not data.startswith("@"):
        return None
    return data


def append_bib_entry(root: Path, entry: str, cite_key: str):
    """Append a BibTeX entry to .crucible/references.bib if not already present.

    Uses file locking and atomic writes to prevent concurrent writers
    from clobbering each other's entries.
    """
    bib_path = root / CRUCIBLE_DIR / "references.bib"
    lock_path = bib_path.with_suffix(".bib.lock")

    with FileLock(lock_path, timeout=10):
        if bib_path.exists():
            existing = bib_path.read_text(encoding="utf-8")
            if f"{{{cite_key}," in existing:
                return
            new_content = existing.rstrip() + "\n\n" + entry + "\n"
        else:
            new_content = entry + "\n"
        atomic_write_text(bib_path, new_content)


def ingest_source(
    root: Path,
    db: CrucibleDB,
    source_path: Path,
    title: str | None = None,
    source_type: str | None = None,
    shareable: bool | None = None,
    url: str | None = None,
    date: str | None = None,
    authors: list[str] | None = None,
    doi: str | None = None,
    bibtex: str | None = None,
) -> dict:
    """Ingest a source file into Crucible.

    Returns a dict with source_id, stored_path, extracted_text_path, cite_key, and metadata.
    """
    source_path = source_path.resolve()
    if not source_path.exists():
        raise FileNotFoundError(f"Source not found: {source_path}")

    # Detect type if not specified
    if source_type is None:
        source_type = detect_source_type(source_path)

    # Determine shareability
    if shareable is None:
        shareable = source_type in ("notebook", "data")

    # Copy to destination
    dest_dir = destination_dir(root, source_type, shareable)
    dest_dir.mkdir(parents=True, exist_ok=True)
    dest_path = dest_dir / source_path.name

    # Handle name collisions
    if dest_path.exists() and not dest_path.samefile(source_path):
        stem = source_path.stem
        suffix = source_path.suffix
        counter = 1
        while dest_path.exists():
            dest_path = dest_dir / f"{stem}_{counter}{suffix}"
            counter += 1

    if not dest_path.exists():
        shutil.copy2(str(source_path), str(dest_path))

    # Extract text
    text = extract_text(dest_path)

    # Save extracted text alongside source
    text_path = dest_path.with_suffix(".txt")
    if text:
        atomic_write_text(text_path, text)

    # Determine title
    if title is None:
        title = extract_title(dest_path, text)

    # Use today if no date given
    if date is None:
        date = datetime.now().strftime("%Y-%m-%d")

    # Relative path for database
    rel_path = str(dest_path.relative_to(root))

    # Generate cite key and bib entry. Preference order:
    #   1. user-provided --bibtex string
    #   2. fetched from doi.org content negotiation
    #   3. minimal entry from known metadata
    cite_key = generate_cite_key(title, authors, date)
    bib_entry = None
    if bibtex and bibtex.strip().startswith("@"):
        bib_entry = replace_cite_key(bibtex.strip(), cite_key)
    elif doi:
        fetched = fetch_bibtex_from_doi(doi)
        if fetched:
            bib_entry = replace_cite_key(fetched, cite_key)
    if bib_entry is None:
        bib_entry = generate_bib_entry(
            cite_key, title, authors, date, url, source_type, doi=doi,
        )
    append_bib_entry(root, bib_entry, cite_key)

    # Register in database (store cite_key in metadata)
    source_id = db.add_source(
        path=rel_path,
        title=title,
        source_type=source_type,
        shareable=shareable,
        url=url,
        authors=authors,
        date=date,
        metadata={"cite_key": cite_key, "doi": doi} if doi else {"cite_key": cite_key},
    )

    return {
        "source_id": source_id,
        "title": title,
        "cite_key": cite_key,
        "source_type": source_type,
        "shareable": shareable,
        "stored_path": str(dest_path),
        "relative_path": rel_path,
        "text_path": str(text_path) if text else None,
        "text_length": len(text),
        "date": date,
    }
